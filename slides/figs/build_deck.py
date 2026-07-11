#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build presentation.pptx from the Greek beamer content, styled to match Τελική_Παρουσίαση.pptx
(NTUA navy/gold Calibri template). Native/editable text + tables; equations & result grids embedded
as PNGs. Adds an attention-results slide (baseline vs morph-attn vs linear-attn, all available folds)."""
import json, glob, os, statistics as st
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT   = "/home/Kasimatis/Documents/kasimat/morph-unet"
FIGS   = f"{ROOT}/slides/figs"
RES    = f"{ROOT}/results"
OUTPPT = f"{ROOT}/slides/presentation.pptx"

# ---- design tokens (from Τελική_Παρουσίαση.pptx) ----
NAVY   = RGBColor(0x0F,0x24,0x4F)
NAVYD  = RGBColor(0x08,0x18,0x36)
GOLD   = RGBColor(0xC9,0xA2,0x27)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LBLUE  = RGBColor(0xC0,0xCB,0xE3)
BODY   = RGBColor(0x4A,0x4A,0x4A)
DIV    = RGBColor(0xD8,0xDC,0xE3)
GOODBG = RGBColor(0xE2,0xF0,0xEA)
FONT   = "Calibri"
FOOTER = "ΕΜΠ  ·  Σχολή ΗΜΜΥ  ·  Όραση Υπολογιστών  ·  Compact Morphological/Linear U-Net"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
EMU = 914400

def _norm(v): return int(v) if isinstance(v,(int,float)) and v>1000 else Inches(v)

def rect(slide, l, t, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p

def run(p, text, size, color, bold=False, italic=False, font=FONT):
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return r

def bullet(p, indent=0.30):
    """turn paragraph into a hanging-indent bullet."""
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(indent))))
    pPr.set('indent', str(int(-Inches(indent))))

def fit(img, box_l, box_t, box_w, box_h):
    iw, ih = Image.open(img).size
    ar = iw/ih; bar = box_w/box_h
    if ar > bar: w = box_w; h = w/ar
    else:        h = box_h; w = h*ar
    return box_l+(box_w-w)/2, box_t+(box_h-h)/2, w, h

# ---------- slide chrome ----------
PAGE = {"i": 1}
TOTAL = None  # set after we know content count

def content_slide(title):
    s = prs.slides.add_slide(BLANK)
    PAGE["i"] += 1
    rect(s, 0, 0, 13.333, 0.82, NAVY)               # top band
    rect(s, 0, 0.82, 13.333, 0.04, GOLD)            # gold rule
    tf = textbox(s, 0.55, 0.0, 12.2, 0.82, MSO_ANCHOR.MIDDLE)
    run(para(tf, True), title, 26, WHITE, bold=True)
    rect(s, 0.55, 7.08, 12.23, 0.012, DIV)          # footer divider
    ftf = textbox(s, 0.55, 7.13, 9.6, 0.30)
    run(para(ftf, True), FOOTER, 9.5, BODY)
    ptf = textbox(s, 11.0, 7.12, 1.78, 0.30)
    pp = para(ptf, True); pp.alignment = PP_ALIGN.RIGHT
    run(pp, f"{PAGE['i']} / {TOTAL}", 11, NAVY, bold=True)
    return s

def body_box(s, top=1.12, height=5.85):   # standard body text region
    return textbox(s, 0.55, top, 12.23, height)

def add_bullets(tf, items, size=18, first=True, gap=8):
    for k, it in enumerate(items):
        p = para(tf, first and k == 0)
        p.space_after = Pt(gap); p.line_spacing = 1.02
        run(p, "•  ", size, GOLD, bold=True)
        run(p, it, size, BODY)
        bullet(p)

def intro(tf, text, size=19, first=True, gap=10):
    p = para(tf, first); p.space_after = Pt(gap); p.line_spacing = 1.02
    run(p, text, size, BODY)

def image_slide(title, img, caption=None):
    s = content_slide(title)
    if caption:
        ctf = textbox(s, 0.55, 1.10, 12.23, 0.66)
        run(para(ctf, True), caption, 13.5, BODY)
        box = (0.55, 1.86, 12.23, 5.05)
    else:
        box = (0.55, 1.15, 12.23, 5.78)
    l, t, w, h = fit(f"{FIGS}/{img}", *box)
    s.shapes.add_picture(f"{FIGS}/{img}", Inches(l), Inches(t), Inches(w), Inches(h))
    return s

def emb(s, img, l, t, w):   # place image at width w (in), top t; returns bottom
    iw, ih = Image.open(f"{FIGS}/{img}").size
    h = w*ih/iw
    s.shapes.add_picture(f"{FIGS}/{img}", Inches(l), Inches(t), Inches(w), Inches(h))
    return t+h

# ---------- table helper ----------
def set_table_plain(table):
    tbl = table._tbl; tblPr = tbl.tblPr
    for e in tblPr.findall(qn('a:tableStyleId')): tblPr.remove(e)
    sid = tblPr.makeelement(qn('a:tableStyleId'), {}); sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tblPr.append(sid)
    table.first_row = False; table.horz_banding = False

def cell(c, text, size, color, bold=False, fill=None, align=PP_ALIGN.LEFT):
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    else:
        c.fill.background()
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    c.margin_left = Inches(0.10); c.margin_right = Inches(0.08)
    c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
    tf = c.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run(p, text, size, color, bold=bold)

def make_table(s, left, top, colw, rows, header, rowh=0.42, hsize=13.5, bsize=13.5):
    nR = len(rows)+1; nC = len(colw)
    gf = s.shapes.add_table(nR, nC, Inches(left), Inches(top), Inches(sum(colw)), Inches(rowh*nR))
    t = gf.table; set_table_plain(t)
    for j, w in enumerate(colw): t.columns[j].width = Inches(w)
    for i in range(nR): t.rows[i].height = Inches(rowh)
    for j, htxt in enumerate(header):
        cell(t.cell(0, j), htxt, hsize, WHITE, bold=True, fill=NAVY,
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, rowvals in enumerate(rows, start=1):
        for j, val in enumerate(rowvals):
            cell(t.cell(i, j), val, bsize, NAVY if j == 0 else BODY,
                 bold=(j == 0), fill=WHITE,
                 align=PP_ALIGN.LEFT if j <= 1 else PP_ALIGN.CENTER)
    # gold underline under header + light row separators
    rect(s, left, top+rowh, sum(colw), 0.028, GOLD)
    for i in range(2, nR):
        rect(s, left, top+rowh*i, sum(colw), 0.010, DIV)
    return gf

# ---------- attention results (computed from JSONs) ----------
def stat(tag):
    dv, dt, av, at = [], [], [], []
    for f in sorted(glob.glob(f"{RES}/{tag}_f*_scores.json")):
        r = json.load(open(f))["results"]["mean"]; v, t = r["Vessel"], r["Tumour"]
        dv.append(v["Dice"]); dt.append(t["Dice"])
        av.append(v["Avg. Symmetric Surface Distance"]); at.append(t["Avg. Symmetric Surface Distance"])
    n = len(dv)
    def f(xs, d=3):
        m = sum(xs)/len(xs); s = st.pstdev(xs) if len(xs) > 1 else 0
        return f"{m:.{d}f}" + (f"±{s:.{d}f}" if s > 0 else "")
    return n, f(dv), f(dt), f(av, 2), f(at, 2)

# =========================================================================
#  TITLE SLIDE
# =========================================================================
def title_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 9.35, 0, 3.983, 7.5, NAVYD)
    rect(s, 9.30, 0, 0.05, 7.5, GOLD)
    # institution
    itf = textbox(s, 0.60, 0.55, 8.4, 1.3)
    run(para(itf, True), "ΕΘΝΙΚΟ ΜΕΤΣΟΒΙΟ ΠΟΛΥΤΕΧΝΕΙΟ", 15, WHITE, bold=True)
    p = para(itf); run(p, "Σχολή Ηλεκτρολόγων Μηχανικών και Μηχανικών Υπολογιστών", 12.5, LBLUE)
    p = para(itf); run(p, "Μάθημα: Όραση Υπολογιστών", 12.5, LBLUE)
    rect(s, 0.62, 2.75, 0.90, 0.07, GOLD)
    # title
    ttf = textbox(s, 0.60, 3.02, 8.55, 1.9)
    p = para(ttf, True); p.line_spacing = 1.03
    run(p, "Ένα συμπαγές", 36, WHITE, bold=True)
    p = para(ttf); p.line_spacing = 1.03
    run(p, "Morphological / Linear U-Net block", 36, WHITE, bold=True)
    # subtitle (gold)
    stf = textbox(s, 0.60, 5.25, 8.55, 1.0)
    p = para(stf, True); run(p, "Αποδοτικό segmentation και channel pruning", 18, GOLD, bold=True)
    p = para(stf); run(p, "στο MSD Task08 — hepatic vessel", 15, LBLUE)
    # date
    dtf = textbox(s, 0.60, 6.78, 8.0, 0.4)
    run(para(dtf, True), "Αθήνα, ακαδημαϊκό έτος 2025–2026", 11.5, LBLUE)
    # team panel
    gtf = textbox(s, 9.62, 2.80, 3.4, 3.8)
    run(para(gtf, True), "ΟΜΑΔΑ", 12, GOLD, bold=True)
    for name in ["Μπακογιώργος Ηλίας", "Κασιμάτης Δημήτριος", "Παναγόπουλος Μιχαήλ",
                 "Ρήγας Αθανάσιος", "Φειδάκης Αθανάσιος"]:
        p = para(gtf); p.space_before = Pt(6); run(p, name, 12.5, WHITE)

# =========================================================================
#  CONTENT SLIDES  (list of builder callables -> lets us count TOTAL)
# =========================================================================
def s_building():
    s = content_slide("Πειράματα με Morphology και το building block")
    tf = body_box(s, 1.12, 3.1)
    add_bullets(tf, [
        "Η max/min-plus (tropical) άλγεβρα γενικεύει το μη γραμμικό μέρος ενός CNN. Ένα max-pool δεν είναι "
        "παρά ένα max-plus dilation, οπότε ένα morphological layer περιέχει ήδη το pooling και πάει παραπέρα, "
        "παίρνοντας το max ή το min πάνω σε ένα εκπαιδευόμενο structuring element αντί για ένα σταθερό παράθυρο.",
        "Περιορισμένος χρόνος και hardware μάς οδήγησαν σε ένα compact block που σπάει τη συνέλιξη σε ένα "
        "3×3 depthwise filter (κάθε map βλέπει μόνο το δικό του input channel, χωρίς mixing) και ένα 1×1 "
        "linear mixer που αναμειγνύει τα channels και κρατά το residual και το ReLU/pooling.",
    ], size=18)
    b = emb(s, "eq_params.png", 1.9, 4.55, 9.5)
    ntf = textbox(s, 0.55, b+0.10, 12.23, 0.4); p = para(ntf, True); p.alignment = PP_ALIGN.CENTER
    run(p, "C₍ᵢₙ₎: channels εισόδου · C₍ₒᵤₜ₎: channels εξόδου · η 3×3 convolution είναι depthwise", 12, BODY)

def s_neuron():
    s = content_slide("Ο morphological νευρώνας (MPM / R-MPM)")
    tf = body_box(s, 1.12, 0.8)
    intro(tf, "Ένας MPM νευρώνας συνδυάζει ένα max-plus dilation και ένα min-plus erosion που μοιράζονται τα "
              "ίδια βάρη w (το structuring element) και χρησιμοποιούν ξεχωριστά biases.", 18)
    emb(s, "eq_mpm.png", 0.9, 1.95, 11.5)
    tf2 = body_box(s, 3.05, 3.9)
    add_bullets(tf2, [
        "Επειδή το dilation και το erosion μοιράζονται βάρη, η έξοδος έχει σχεδόν μηδενική μέση τιμή· ένα "
        "linear 1×1 unit στη συνέχεια αναμειγνύει τα channels ρυθμίζοντας τη διακύμανση και δημιουργώντας "
        "έναν universal approximator.",
        "Με Maslov dequantisation μετατρέπουμε το αραιό max σε ομαλό και εκπαιδεύσιμο:   "
        "maxₖ aₖ ≈ (1/β)·log Σₖ exp(β·aₖ).",
        "Το R-MPM προσθέτει ένα residual skip γύρω από το block, y = x + block(x), ώστε ένα μόνο unit να "
        "γενικεύει καλύτερα και να εκπαιδεύεται ευκολότερα. Επιπλέον τα morphological layers εξάγουν γεωμετρικά "
        "χαρακτηριστικά που χρειάζονται λεπτή δομή, κάτι που επιτυγχάνει το residual skip.",
    ], size=17)

def s_models():
    s = content_slide("Πειράματα: τα μοντέλα")
    tf = body_box(s, 1.12, 2.0)
    add_bullets(tf, [
        "Το baseline είναι ένα residual U-Net, για καλύτερη σύγκριση.",
        "Το Compact-CNN χρησιμοποιεί ένα linear 3×3, που κάνει το block μια depthwise-separable συνέλιξη.",
        "Το Morphological-CNN χρησιμοποιεί ένα morphological 3×3 (χωρίς channel mixing) σε κάποια layers και "
        "κρατά standard linear blocks στα υπόλοιπα.",
    ], size=18)
    make_table(s, 1.2, 3.35, [2.3, 5.3, 3.3],
               [["bottleneck", "center", "μόνο το βαθύτερο"],
                ["deep", "enc4, center, dec4", "βαθιά"],
                ["full_l2", "enc3, enc4, center, dec4, dec3", "μεσαία και βαθιά"],
                ["heavy", "και τα εννέα stages", "παντού"],
                ["balanced", "enc1, enc2, dec1, dec2", "μόνο υψηλή ανάλυση"]],
               ["config", "morphological stages", "ανάλυση"], rowh=0.52)

def s_dice():   image_slide("Αποτελέσματα Dice / ASSD (test set)", "tbl_slide7_results.png")

def s_arch_concl():
    s = content_slide("Συμπεράσματα για την αρχιτεκτονική")
    tf = body_box(s)
    add_bullets(tf, [
        "Το linear Compact-CNN ξεπερνά και τα δύο baselines των U-Nets χρησιμοποιώντας περίπου το ένα πέμπτο "
        "των parameters.",
        "Το U-Net με morphological blocks επιτυγχάνει επιδόσεις κοντά στο linear Compact-CNN· τα καλύτερα "
        "morphological configurations έχουν τα morphology blocks στο κέντρο, κάτι που συμφωνεί με άλλες αναφορές "
        "ότι η morphology έχει ελαφρώς λιγότερη ακρίβεια.",
        "Τα morphology blocks δεν αποδίδουν καλά στα υψηλότερα layers. Οι max/min operators κρατούν μόνο την "
        "ισχυρότερη απόκριση σε μια γειτονιά· στην υψηλή ανάλυση το σήμα είναι γεμάτο σημαντικές λεπτομέρειες, "
        "οπότε χάνεται γεωμετρία που τα skip connections δεν μπορούν πια να ανακτήσουν (και αντιδρά άσχημα σε "
        "salt-and-pepper ακραίες τιμές).",
        "Στο bottleneck η ίδια αραιή επιλογή δεν χάνει σχεδόν τίποτα χάρη στα skip connections, και μπορεί "
        "ακόμη να προσθέσει σημαντική πληροφορία από ένα κυρίαρχο pixel της γειτονιάς.",
    ], size=18)

def s_prune_intro():
    s = content_slide("Pruning του compact block")
    tf = body_box(s)
    intro(tf, "Οι 1×1 συνελίξεις δίνουν ένα βάρος σε κάθε channel, και τις χρησιμοποιούμε για να κάνουμε "
              "prune τα αχρησιμοποίητα channels.", 19)
    add_bullets(tf, [
        "Συγκρίνουμε έξι scoring schemes — l1x1, lin, act και morph, μαζί με το fb και το προτεινόμενο fb-new "
        "— έναντι ενός random baseline.",
        "Τα channels που κρατάμε κατανέμονται με δύο τρόπους: το local κρατά ένα σταθερό ratio σε κάθε layer, "
        "ενώ το global ξοδεύει ένα ενιαίο budget σε όλο το μοντέλο και κρατά channels άνισα.",
    ], size=18, first=False)

def s_schemes():
    s = content_slide("Τα τέσσερα βασικά scoring schemes")
    tf = body_box(s, 1.12, 0.55)
    intro(tf, "Κάθε scheme δίνει σε κάθε input channel i ένα score, και κρατάμε αυτά με το υψηλότερο score.", 18)
    emb(s, "schemes_table.png", 0.7, 1.75, 11.9)
    tf2 = body_box(s, 4.95, 2.4)
    add_bullets(tf2, [
        "‖proj‖ — πόσο έντονα οι 1×1 mixers διαβάζουν το κάθε input channel.",
        "|α| — το gain που έχει μάθει το unit για κάθε pixel του 3×3 kernel.",
        "𝔼|morph(x)| — η μέση απόλυτη τιμή των activations όλων των pixels του channel.",
        "spread(SE) = max − min των βαρών του Structuring Element: πόσο αναδιαμορφώνει τη γειτονιά "
        "(≈0 → απλό max-pooling).",
    ], size=14, gap=4)

def s_fb():
    s = content_slide("Το scheme fb και το fb-new")
    tf = body_box(s, 1.12, 0.75)
    intro(tf, "Το fb αντιμετωπίζει την αλυσίδα των morph units ως HMM και βαθμολογεί κάθε channel με βάση τη "
              "global posterior occupancy του  γ = α·β.", 18)
    tf2 = body_box(s, 1.95, 2.75)
    add_bullets(tf2, [
        "Το unigram prior είναι π ∝ 𝔼|morph|, και η transition T(i→j) είναι η co-activation μεταξύ διαδοχικών "
        "layers, κανονικοποιημένη ανά γραμμή με τη ℓ₁ norm.",
        "Ένα forward pass α μετρά πόσο προσβάσιμο είναι ένα channel από τα προηγούμενα, κι ένα backward pass β "
        "πόσο επηρεάζει τα επόμενα.",
        "Το global allocation ταιριάζει καλά στο fb, καθώς αποτυπώνει τη συνολική πιθανότητα των συνδυασμένων "
        "activations.",
    ], size=17)
    # fb-new highlighted block
    rect(s, 0.55, 5.35, 12.23, 1.45, GOODBG)
    rect(s, 0.55, 5.35, 0.10, 1.45, GOLD)
    btf = textbox(s, 0.85, 5.48, 11.75, 1.2)
    p = para(btf, True); run(p, "fb-new  (προτεινόμενο, μελλοντική εργασία)", 15, NAVY, bold=True)
    p = para(btf); p.space_before = Pt(3)
    run(p, "Περιορίζει τα στατιστικά των activations στα foreground pixels (vessel & tumour) χρησιμοποιώντας το "
           "receptive field κάθε channel — τα decoder layers κληρονομούν το receptive field του απέναντι encoder "
           "layer, ώστε να αγνοούμε τα background pixels.", 13.5, BODY)

def _prune(title, img, cap): return lambda: image_slide(title, img, cap)
def _chan(title, img):       return lambda: image_slide(title, img)

def s_prune_concl():
    s = content_slide("Συμπεράσματα για το pruning")
    tf = body_box(s, 1.12, 2.7)
    add_bullets(tf, [
        "Το pruning συμπεριφέρεται σχεδόν ίδια σε όλα τα μοντέλα, αφού βασίζεται κυρίως στο 1×1 separable "
        "convolution. Τα center layers σβήνονται πιο εύκολα, επειδή έχουν το περισσότερο redundancy — μεγάλο "
        "μέρος της πληροφορίας μεταφέρεται μέσω των skip connections.",
        "Το fb τείνει να κρατά ισχυρά συσχετισμένα paths, οπότε διατηρεί πλεοναστικούς γείτονες και πετά "
        "χρήσιμα channels χαμηλής πιθανότητας (το fb-new διορθώνει ακριβώς αυτό).",
    ], size=18)
    ntf = textbox(s, 0.55, 4.15, 12.23, 2.7)
    p = para(ntf, True); p.line_spacing = 1.05
    run(p, "Το linear Compact-CNN block κάνει prune τουλάχιστον εξίσου καλά με τα morphological και παραμένει "
           "πιο ακριβές ακόμη κι όταν το κάνουμε prune στο έπακρο· είναι δε ένα linear building block για το "
           "οποίο υπάρχουν έτοιμες βιβλιοθήκες CUDA και hardware accelerators. Τα πειράματα δείχνουν ότι τα "
           "υπάρχοντα morphological layers δεν βελτιώνουν ιδιαίτερα το pruning ή την ακρίβεια. Ωστόσο, η "
           "υλοποίησή μας δεν εκμεταλλεύεται πλήρως τη δομή pruning που προτείνεται για MNNs με πλήρη "
           "3×3×C₍ᵢₙ₎ morphological kernels (αναγκαία σχεδιαστική επιλογή λόγω υπολογιστικού κόστους).", 15, BODY)

def s_attn_discuss():
    s = content_slide("Morphological vs. linear attention στο Res-U-Net")
    tf = body_box(s)
    add_bullets(tf, [
        "Το morphological attention χρησιμοποιεί top-hat και bottom-hat φίλτρα στην είσοδο, ενώνει το "
        "αποτέλεσμα κι εφαρμόζει ένα linear gate από πάνω. Ενσωματώνει ένα ισχυρό γεωμετρικό prior για τα "
        "λεπτά, φωτεινά hepatic vessels.",
        "Το linear attention καλείται να μάθει αυτούς τους γεωμετρικούς operators από τα δεδομένα, που εδώ "
        "είναι λιγοστά· έτσι το top-hat / bottom-hat prior του morphological attention αναμένεται να δώσει "
        "καλύτερη επίδοση.",
        "Η residual σύνδεση μεταφέρει λεπτομέρεια κατά μήκος του identity path της, ώστε τα skip connections "
        "να φέρνουν πιο λεπτομερείς αναπαραστάσεις στον decoder.",
        "Υβριδική προσέγγιση: morphological attention στα δύο πρώτα, χαμηλού επιπέδου skip levels και linear ή "
        "καθόλου attention στα βαθύτερα, πιο σημασιολογικά επίπεδα.",
    ], size=18)

def s_attn_results():
    s = content_slide("Αποτελέσματα attention (test set)")
    tf = body_box(s, 1.12, 0.5)
    intro(tf, "Τρία matched arms στο ίδιο residual U-Net backbone. Μέση τιμή ± τυπική απόκλιση στα διαθέσιμα "
              "folds.", 17)
    rows = []
    for tag, label in [("unet_baseline", "Baseline residual U-Net"),
                       ("unet_morphattn", "Morphological attention"),
                       ("unet_linattn_g05", "Linear attention (skip γ₀=0.5)")]:
        n, dv, dt, av, at = stat(tag)
        rows.append([label, str(n), dv, dt, av, at])
    make_table(s, 0.75, 1.95, [3.55, 1.0, 2.15, 2.15, 1.85, 1.85], rows,
               ["μοντέλο", "folds", "Dice Vessel", "Dice Tumour", "ASSD Vessel ↓", "ASSD Tumour ↓"],
               rowh=0.62, hsize=13, bsize=14)
    ntf = textbox(s, 0.55, 5.15, 12.23, 1.7)
    p = para(ntf, True); p.line_spacing = 1.05
    run(p, "Το linear attention εδώ ξεκινά το skip gate μισάνοιχτο (ReZero γ₀=0.5) αντί για 0. Στο fold 0 "
           "αυτό βελτίωσε αισθητά και το Vessel και το Tumour Dice έναντι του κλασικού (γ₀=0). Baseline 3 "
           "folds, morphological attention 2, linear attention 1 fold προς το παρόν (τα folds 1–2 στην ουρά "
           "μετά το pruning). Το Tumour είναι σπάνιο στο Task08, οπότε οι μετρικές του έχουν μεγάλη διακύμανση.",
           13.5, BODY)

def s_further():
    s = content_slide("Περαιτέρω πειραματισμός")
    tf = body_box(s)
    add_bullets(tf, [
        "Δοκιμή του block σε περισσότερα datasets και άλλα objectives, όπως classification εκτός από segmentation.",
        "Συμπλήρωση των cross-fold pruning και των Dice/ASSD μετρικών σε κάθε περίπτωση, καθώς αρκετά runs ήταν "
        "single-fold λόγω υπολογιστικού κόστους.",
        "Δοκιμή του fb-new σε περισσότερα datasets/setups με το compact block. Θα μπορούσε να εξοικονομήσει "
        "πολλά parameters σε block-based networks, π.χ. το MobileNet, που χρησιμοποιεί ένα πολύ παρόμοιο "
        "depthwise-separable block.",
    ], size=18)

def s_refs():
    s = content_slide("Αναφορές")
    tf = body_box(s)
    add_bullets(tf, [
        "P. Maragos — Morphological systems / max-plus operators.",
        "Maragos & Fotopoulos — MPM / R-MPM neuron.",
        "Depthwise-separable blocks (MobileNet).",
        "NISP-style importance propagation (Taylor).",
    ], size=18)

CONTENT = [
    s_building, s_neuron, s_models, s_dice, s_arch_concl, s_prune_intro, s_schemes, s_fb,
    _prune("Αποτελέσματα pruning για το mpm_deep", "prune_deep.png",
           "Στο deep, η morphology βρίσκεται στα τρία stages γύρω από το bottleneck (enc4, center, dec4)· το υπόλοιπο U-Net χρησιμοποιεί standard blocks."),
    _prune("Αποτελέσματα pruning για το mpm_bottleneck", "prune_bottleneck.png",
           "Στο bottleneck, μόνο το βαθύτερο stage (center) είναι morphological. Είναι το μεγαλύτερο morphological μοντέλο (18.5M)."),
    _prune("Αποτελέσματα pruning για το mpm_full_l2", "prune_full_l2.png",
           "Στο full_l2, η morphology καλύπτει το βαθύτερο μισό (enc3, enc4, center, dec4, dec3)· τα δύο επίπεδα υψηλότερης ανάλυσης παραμένουν linear. Το ισχυρότερο morphological μοντέλο."),
    _prune("Αποτελέσματα pruning για το morphunet_heavy", "prune_heavy.png",
           "Στο heavy, καθένα από τα εννέα stages είναι morphological. Με 6M είναι από τα μικρότερα μοντέλα, καθώς η morphology είναι φθηνότερη στα layers με πολλά channels."),
    _prune("Αποτελέσματα pruning για το linear Compact-CNN", "prune_convsep.png",
           "Το linear δίδυμο του heavy: ίδια διάταξη block παντού αλλά με depthwise linear filters αντί για morphology (κανένα morphological layer)."),
    _chan("Κατανομή channels του deep", "chan_deep.png"),
    _chan("Κατανομή channels του bottleneck", "chan_bottleneck.png"),
    _chan("Κατανομή channels του morphological (heavy)", "chan_heavy.png"),
    _chan("Κατανομή channels του linear (heavy)", "chan_linear.png"),
    _chan("Κατανομή channels του full_l2", "chan_full_l2.png"),
    s_prune_concl,
    _chan("Τα βαθύτερα layers κλαδεύονται πολύ πιο έντονα", "plot_c_retention_by_resolution.png"),
    s_attn_discuss, s_attn_results, s_further, s_refs,
]

# extra fold-0 convsep prune slide, added ONLY once its PNG exists (kept alongside the fold-1 slide)
if os.path.exists(f"{FIGS}/prune_convsep_f0.png"):
    _cap0 = ("Ίδιο pruning στο fold 0 (unpruned macro 0.451) αντί για το ασθενέστερο fold 1 (0.368) — "
             "faithful σύγκριση με τα άλλα single-fold μοντέλα (full_l2, heavy). Το fold 1 είχε πολύ "
             "χαμηλό Tumour Dice (0.253), που ρίχνει το macro.")
    CONTENT.insert(13, _prune("Αποτελέσματα pruning: linear Compact-CNN — fold 0", "prune_convsep_f0.png", _cap0))

TOTAL = len(CONTENT) + 1
title_slide()
for fn in CONTENT:
    fn()

prs.save(OUTPPT)
print(f"saved {OUTPPT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, TOTAL={TOTAL})")
